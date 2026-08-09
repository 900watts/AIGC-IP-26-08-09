"""
Build pitch deck for ACTIF A7 - 莞仔的漫博一日
7 slides, soft pastel toy editorial palette (no purple/violet/indigo)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# Palette (locked from DESIGN.md — matches H5)
INK     = RGBColor(0x1A, 0x1F, 0x2E)
CREAM   = RGBColor(0xFF, 0xF5, 0xE6)
CREAM2  = RGBColor(0xFC, 0xE9, 0xCF)
SKY     = RGBColor(0x7F, 0xB8, 0xE8)
SKY_DP  = RGBColor(0x5E, 0x9C, 0xD2)
APRICOT = RGBColor(0xFF, 0x8A, 0x4C)
LONGAN  = RGBColor(0xC9, 0xA3, 0x6A)
PETAL   = RGBColor(0xFF, 0xD9, 0xC2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SHADOW  = RGBColor(0xE5, 0xC9, 0xA0)

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# --------- helpers ---------
def add_rect(slide, x, y, w, h, fill, line=None, shadow=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if False else MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    if not shadow:
        # remove default shadow
        sp = s.shadow
        # python-pptx default has shadow; clear via xml
        spPr = s._element.spPr
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        eff = etree.SubElement(spPr, qn('a:effectLst'))
    return s

def add_round(slide, x, y, w, h, fill, radius_ratio=0.3, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # adjust corner radius
    try:
        s.adjustments[0] = radius_ratio
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    # remove shadow
    spPr = s._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    etree.SubElement(spPr, qn('a:effectLst'))
    return s

def add_text(slide, x, y, w, h, text, size=20, bold=False, color=INK, font='Plus Jakarta Sans', align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_emoji_or_icon(slide, x, y, w, h, char='☆', color=APRICOT, size=40):
    """Use unicode shapes instead of emoji to avoid platform font issues."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = char
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = 'Plus Jakarta Sans'
    return tb

def set_bg(slide, color):
    bg = add_rect(slide, 0, 0, SW, SH, color)
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)

# ---------- SLIDE 1: cover ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, CREAM)
# Decorative orbit rings
add_emoji_or_icon(s, Inches(0.4), Inches(0.4), Inches(0.6), Inches(0.6), char='◎', color=APRICOT, size=28)
# big title
add_text(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2), '莞仔的漫博一日', size=72, bold=True, color=INK, font='Fraunces')
add_text(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.7), 'Guànzǎi · One Day at ACTIF', size=28, color=APRICOT, font='Fraunces')
# accent bar
add_rect(s, Inches(0.8), Inches(3.95), Inches(2.2), Inches(0.06), APRICOT)
# subtitle
add_text(s, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.6), '一个会跟你聊天的 IP 吉祥物 · H5 互动体验', size=22, color=INK)
# bottom meta
add_text(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4), '第十六届中国国际动漫博览会 · AIGC 黑客松', size=16, color=INK)
add_text(s, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.4), 'Group A7 · 2026.08.09 · 东莞石排', size=14, color=LONGAN)
# mascot placeholder (right side) — use sprite image
import os
sprite_path = os.path.join(os.path.dirname(__file__), 'assets', 'mascot-wave.png')
if os.path.exists(sprite_path):
    s.shapes.add_picture(sprite_path, Inches(9.5), Inches(1.4), height=Inches(4.5))

# ---------- SLIDE 2: problem & insight ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '01 · 洞察', size=18, bold=True, color=APRICOT, font='Plus Jakarta Sans')
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '吉祥物不能只是 POSE, 他要会说话', size=44, bold=True, color=INK, font='Fraunces')
# two columns
add_text(s, Inches(0.8), Inches(2.5), Inches(5.7), Inches(0.5), '现在的漫博 IP 传播', size=18, bold=True, color=SKY_DP)
add_text(s, Inches(0.8), Inches(3.0), Inches(5.7), Inches(3.5),
         ['• 海报、贴纸、表情包都是单向输出', '• 用户只能看, 不能跟 IP 互动', '• 同一个 IP, 永远只有一套话', '• 工作人员要回答几千条重复问题', '• 线下流量无法变成可复用的内容资产'],
         size=18, color=INK)
add_text(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.5), '我们想做的', size=18, bold=True, color=APRICOT)
add_text(s, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.5),
         ['• IP 角色可以像真人一样聊', '• 永远在漫博中心门口 "在线"', '• 朋友圈文案、活动推荐、冷知识随手问', '• 用户拿到的不只是看, 是 "我跟 莞仔 聊了"', '• 一次开发, 持续产生新内容'],
         size=18, color=INK)

# ---------- SLIDE 3: solution ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '02 · 方案', size=18, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '把 莞仔 变成可对话的 IP 角色', size=44, bold=True, color=INK, font='Fraunces')
add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(0.6), '一个 H5 落地页, 三件事: 5 站视觉海报 + 莞仔 角色 + 实时聊天', size=20, color=INK)
# three feature cards
def feature_card(slide, x, y, w, h, num, title, body, color):
    add_round(slide, x, y, w, h, WHITE, radius_ratio=0.05)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.6), Inches(0.5), num, size=22, bold=True, color=color)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), w - Inches(0.6), Inches(0.5), title, size=22, bold=True, color=INK, font='Fraunces')
    add_text(slide, x + Inches(0.3), y + Inches(1.3), w - Inches(0.6), h - Inches(1.4), body, size=15, color=INK)

feature_card(s, Inches(0.8), Inches(3.3), Inches(3.9), Inches(3.6), '01', '5 站视觉海报',
             'AI 生成的 莞仔 漫步漫博中心: 清晨·东莞 / 展馆入口 / 潮玩馆 / 签售舞台 / 夜幕·灯笼。每站都藏一个会说话的 IP。',
             APRICOT)
feature_card(s, Inches(4.95), Inches(3.3), Inches(3.9), Inches(3.6), '02', '可点 莞仔',
             '每站都嵌一个会跳 idle/wave/blink 三态的 莞仔。点他一下, 用 AI 模板引擎生成一句带场景感的俏皮话。',
             SKY)
feature_card(s, Inches(9.1), Inches(3.3), Inches(3.9), Inches(3.6), '03', '实时聊天',
             '右下角悬浮气泡, 点开就是完整聊天抽屉。调用本地 Ollama 真实 LLM, 用 莞仔 persona + 漫博知识库回答。',
             LONGAN)

# ---------- SLIDE 4: how it works (architecture) ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '03 · 技术架构', size=18, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '浏览器 ↔ 本地代理 ↔ Ollama', size=44, bold=True, color=INK, font='Fraunces')
# flow boxes
def flow_node(slide, x, y, w, h, title, sub, fill, ink=INK):
    add_round(slide, x, y, w, h, fill, radius_ratio=0.12)
    add_text(slide, x, y + Inches(0.3), w, Inches(0.5), title, size=18, bold=True, color=ink, align=PP_ALIGN.CENTER, font='Fraunces')
    add_text(slide, x, y + Inches(0.85), w, h - Inches(0.9), sub, size=13, color=ink, align=PP_ALIGN.CENTER)

W, H = Inches(2.3), Inches(2.6)
y0 = Inches(2.5)
flow_node(s, Inches(0.8),  y0, W, H, 'H5 前端', 'HTML / CSS / JS\n纯静态\n无构建步骤', CREAM2)
flow_node(s, Inches(3.6),  y0, W, H, '本地代理', 'Python http.server\n/api/ollama\n解决 CORS', SKY, ink=WHITE)
flow_node(s, Inches(6.4),  y0, W, H, 'Ollama', 'glm-5.2:cloud\nqwen3 / deepseek\nfallback 离线', INK, ink=WHITE)
flow_node(s, Inches(9.2),  y0, W, H, '莞仔 角色', 'persona.md\n知识库 + 语气\n≤80 字回复', APRICOT, ink=WHITE)
# arrows
for x in [3.1, 5.9, 8.7]:
    add_text(s, Inches(x), Inches(3.6), Inches(0.4), Inches(0.4), '→', size=36, bold=True, color=LONGAN, align=PP_ALIGN.CENTER)
# footnote
add_text(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5), '关键技术点', size=18, bold=True, color=INK)
add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.4),
         ['• NDJSON 流式响应, token 逐字渲染, 像真人在打字', '• Persona 内嵌完整 ACTIF 知识库, 不会胡说时间地点', '• Ollama 不可用时, 自动 fallback 到 莞仔 口吻的离线模板', '• 同源代理避免浏览器 CORS 阻断, 不依赖任何外部服务'],
         size=15, color=INK)

# ---------- SLIDE 5: AI-created proof ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, CREAM2)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '04 · AIGC 痕迹', size=18, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '每一处都看得到 AI 在干活', size=44, bold=True, color=INK, font='Fraunces')
# 4-cell grid
def aigc_cell(slide, x, y, w, h, title, body, color):
    add_round(slide, x, y, w, h, WHITE, radius_ratio=0.06)
    add_rect(slide, x, y, Inches(0.12), h, color)
    add_text(slide, x + Inches(0.35), y + Inches(0.25), w - Inches(0.5), Inches(0.5), title, size=18, bold=True, color=INK, font='Fraunces')
    add_text(slide, x + Inches(0.35), y + Inches(0.85), w - Inches(0.5), h - Inches(0.95), body, size=14, color=INK)

cw, ch = Inches(5.9), Inches(2.5)
aigc_cell(s, Inches(0.8), Inches(2.4), cw, ch, '7 张图片全部 AI 生成',
          '5 张场景插画 + 1 张 3 态 sprite + 1 张分享卡背景, 用统一 prompt 工程保证 莞仔 不换脸、不换衣。', APRICOT)
aigc_cell(s, Inches(6.95), Inches(2.4), cw, ch, '全站中文文案 AI 起草',
          '海报 5 站标题、Tag、口癖、persona 文档, 全部由我作为语言模型在构建期生成, 不用模板句。', SKY)
aigc_cell(s, Inches(0.8), Inches(5.05), cw, ch, '实时聊天 = 真实 LLM',
          'Ollama glm-5.2:cloud 流式响应, 不是写死的假对话。每条回复都带 thinking 链路, 可审计。', LONGAN)
aigc_cell(s, Inches(6.95), Inches(5.05), cw, ch, '离线 persona 兜底',
          'Ollama 挂了, 自动切到 persona 模板, 莞仔 不会装死, 现场演示永远跑得通。', INK)

# ---------- SLIDE 6: demo ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, WHITE)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '05 · 现场演示', size=18, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '30 秒看懂', size=44, bold=True, color=INK, font='Fraunces')
# 3 demo beats
beat_w = Inches(3.9)
def demo_beat(slide, x, y, w, h, num, title, body, color):
    add_round(slide, x, y, w, h, color, radius_ratio=0.1)
    add_text(slide, x, y + Inches(0.3), w, Inches(0.7), num, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Fraunces')
    add_text(slide, x, y + Inches(1.2), w, Inches(0.5), title, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font='Fraunces')
    add_text(slide, x + Inches(0.3), y + Inches(1.85), w - Inches(0.6), h - Inches(2), body, size=14, color=WHITE, align=PP_ALIGN.CENTER)

demo_beat(s, Inches(0.8),  Inches(2.5), beat_w, Inches(4.3), '01', '打开 H5',
          '5 站海报自动滑入, 顶部计数, 莞仔 sprite 待机 idle 态', SKY)
demo_beat(s, Inches(4.95), Inches(2.5), beat_w, Inches(4.3), '02', '点 莞仔',
          '每站 sprite 切到 wave/blink, AI 模板生成一句带场景的口癖', APRICOT)
demo_beat(s, Inches(9.1),  Inches(2.5), beat_w, Inches(4.3), '03', '聊起来',
          '右下角悬浮气泡 → 聊天抽屉 → Ollama 流式输出 persona 回复', LONGAN)

# ---------- SLIDE 7: impact & next ----------
s = prs.slides.add_slide(BLANK)
set_bg(s, CREAM)
add_text(s, Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.6), '06 · 价值 & 未来', size=18, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(1.1), Inches(11.7), Inches(1.0), '一次开发, 莞仔 永远在', size=44, bold=True, color=INK, font='Fraunces')
# value column
add_text(s, Inches(0.8), Inches(2.5), Inches(5.7), Inches(0.5), '对漫博主办方', size=20, bold=True, color=APRICOT)
add_text(s, Inches(0.8), Inches(3.0), Inches(5.7), Inches(3.5),
         ['• 节省客服人力: 用户 80% 的问题可以交给 莞仔',
          '• 提升停留时长: 互动 > 浏览, 转化率更高',
          '• 内容持续产出: 每次问都是一次新内容',
          '• 数据可回流: 收集用户问什么, 优化下一年展馆'],
         size=16, color=INK)
add_text(s, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.5), '可扩展方向', size=20, bold=True, color=SKY_DP)
add_text(s, Inches(7.0), Inches(3.0), Inches(5.5), Inches(3.5),
         ['• 接企业微信 / 公众号, 让 莞仔 进私域',
          '• 配 TTS 让 莞仔 开口说, 做漫博语音导览',
          '• 多 IP 套壳: 漫妹、鹤先生、DreamStory 都能用',
          '• 训练专属 LoRA, 莞仔 风格 + 主办方知识库'],
         size=16, color=INK)
# final accent
add_round(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), INK, radius_ratio=0.5)
add_text(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), 'THANK YOU · 记得来 ACTIF 找 莞仔 玩', size=22, bold=True, color=CREAM, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Save
out = os.path.join(os.path.dirname(__file__), 'AIGC_A7_莞仔的漫博一日.pptx')
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides))